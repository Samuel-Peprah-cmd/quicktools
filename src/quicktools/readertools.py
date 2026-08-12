import os

def parse_document_for_reader(file_path: str) -> list[dict]:
    """
    Parses a document and returns a list of paginated text blocks.
    Perfectly formatted for a React frontend to read aloud.
    Returns: [{"page": 1, "text": "Hello world"}, ...]
    """
    if not os.path.exists(file_path):
        raise FileNotFoundError(f"Could not find the file: {file_path}")

    ext = os.path.splitext(file_path)[1].lower()
    
    if ext == '.pdf':
        return _parse_pdf(file_path)
    elif ext == '.docx':
        return _parse_docx(file_path)
    elif ext == '.pptx':
        return _parse_pptx(file_path)
    else:
        raise ValueError(f"File type '{ext}' is not supported. Please upload a PDF, DOCX, or PPTX.")

def _parse_pdf(file_path: str) -> list[dict]:
    try:
        import fitz  # PyMuPDF
    except ImportError:
        raise RuntimeError("The 'pymupdf' package is missing. Please install it.")
        
    pages_data = []
    with fitz.open(file_path) as doc:
        for page_num in range(len(doc)):
            page = doc.load_page(page_num)
            text = page.get_text("text").strip()
            
            # Only add pages that actually have text (skip blank pages)
            if text:
                # Clean up weird PDF line breaks to make it sound natural when read aloud
                clean_text = " ".join(text.split())
                pages_data.append({"page": page_num + 1, "text": clean_text})
                
    return pages_data

def _parse_docx(file_path: str) -> list[dict]:
    try:
        import docx
    except ImportError:
        raise RuntimeError("The 'python-docx' package is missing. Please install it.")
        
    doc = docx.Document(file_path)
    pages_data = []
    
    # Word documents don't have "pages" in code, they have paragraphs.
    # We will group every 5 paragraphs into a "page" block for the UI to digest easily.
    current_text = []
    page_counter = 1
    
    for para in doc.paragraphs:
        para_text = para.text.strip()
        if para_text:
            current_text.append(para_text)
            
            if len(current_text) >= 5:
                pages_data.append({"page": page_counter, "text": " ".join(current_text)})
                current_text = []
                page_counter += 1
                
    # Catch any remaining paragraphs at the end of the document
    if current_text:
        pages_data.append({"page": page_counter, "text": " ".join(current_text)})
        
    return pages_data

def _parse_pptx(file_path: str) -> list[dict]:
    try:
        from pptx import Presentation
    except ImportError:
        raise RuntimeError("The 'python-pptx' package is missing. Please install it.")
        
    prs = Presentation(file_path)
    pages_data = []
    
    # Each slide acts as a page
    for i, slide in enumerate(prs.slides):
        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_text.append(shape.text.strip())
                
        if slide_text:
            clean_text = " ".join(" ".join(slide_text).split())
            pages_data.append({"page": i + 1, "text": clean_text})
            
    return pages_data