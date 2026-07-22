"""PowerPoint utilities: reading, creating .pptx files, and content-based conversion
to/from Word documents and PDFs — with styled themes and image placement.

Note on conversions: Word docs, PDFs, and PowerPoint decks use fundamentally different
document models (flowing paragraphs vs. discrete slides). These functions preserve and
reorganize the actual text content and images sensibly into the new format — they do not
attempt to replicate exact visual layout, because no such 1:1 mapping exists between formats.

Themes are applied entirely in code (background/title/accent colors) — no external
template files or network downloads are required.
"""
import io
import random
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE

# Each theme is a self-contained color palette applied entirely in code.
THEMES = [
    {"name": "Ocean Blue", "background": RGBColor(0x0B, 0x2A, 0x4A), "title": RGBColor(0xFF, 0xFF, 0xFF), "accent": RGBColor(0x9C, 0xD3, 0xFF)},
    {"name": "Sunset Orange", "background": RGBColor(0x3A, 0x1A, 0x0A), "title": RGBColor(0xFF, 0xD9, 0xA0), "accent": RGBColor(0xFF, 0xFF, 0xFF)},
    {"name": "Forest Green", "background": RGBColor(0x0E, 0x2E, 0x1D), "title": RGBColor(0xFF, 0xFF, 0xFF), "accent": RGBColor(0xA8, 0xE6, 0xB8)},
    {"name": "Monochrome", "background": RGBColor(0x1A, 0x1A, 0x1A), "title": RGBColor(0xFF, 0xFF, 0xFF), "accent": RGBColor(0xCC, 0xCC, 0xCC)},
    {"name": "Soft Light", "background": RGBColor(0xF5, 0xF5, 0xF0), "title": RGBColor(0x22, 0x22, 0x22), "accent": RGBColor(0x44, 0x44, 0x44)},
]

class Theme:
    """Pre-defined theme constants for easy IDE auto-completion."""
    OCEAN_BLUE = "Ocean Blue"
    SUNSET_ORANGE = "Sunset Orange"
    FOREST_GREEN = "Forest Green"
    MONOCHROME = "Monochrome"
    SOFT_LIGHT = "Soft Light"


def _pick_theme(theme_name: str | None = None) -> dict:
    """Return a theme dict by name, or a random one if no name is given."""
    if theme_name:
        for t in THEMES:
            if t["name"].lower() == theme_name.lower():
                return t
        raise ValueError(f"Unknown theme '{theme_name}'. Available: {[t['name'] for t in THEMES]}")
    return random.choice(THEMES)


def _apply_background(slide, theme: dict) -> None:
    """Fill a slide's background with the theme's background color."""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = theme["background"]


def _add_slide_with_text(prs: Presentation, theme: dict, title: str, body_lines: list[str],
                          image_path: str | None = None) -> None:
    """Internal helper: add a themed slide with a title, bullet text, and an optional image."""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    _apply_background(slide, theme)

    has_image = image_path is not None
    text_width = Inches(4.7) if has_image else Inches(9)

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), text_width, Inches(1))
    tf = title_box.text_frame
    tf.word_wrap = True
    tf.text = title
    tf.paragraphs[0].font.size = Pt(28)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = theme["title"]

    body_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), text_width, Inches(5.3))
    body_tf = body_box.text_frame
    body_tf.word_wrap = True
    for i, line in enumerate(body_lines):
        p = body_tf.paragraphs[0] if i == 0 else body_tf.add_paragraph()
        p.text = f"• {line}"
        p.font.size = Pt(18)
        p.font.color.rgb = theme["accent"]

    if has_image:
        slide.shapes.add_picture(image_path, Inches(5.5), Inches(1.5), width=Inches(4), height=Inches(4.5))


def read_pptx_text(path: str) -> str:
    """Extract all text from every slide in a .pptx file."""
    prs = Presentation(path)
    lines = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = "".join(run.text for run in paragraph.runs)
                    if text:
                        lines.append(text)
    return "\n".join(lines)


def get_pptx_slide_count(path: str) -> int:
    """Return the number of slides in a .pptx file."""
    return len(Presentation(path).slides)


def list_available_themes() -> list[str]:
    """Return the names of all available built-in themes."""
    return [t["name"] for t in THEMES]


def create_pptx_from_text(text: str, output_path: str, theme_name: str | None = None) -> None:
    """Create a themed .pptx file from text, splitting on blank lines into separate slides.
    Each slide's first line becomes its title; remaining lines become bullet content.
    theme_name picks a specific theme (see list_available_themes()); omit for a random one."""
    theme = _pick_theme(theme_name)
    prs = Presentation()
    blocks = [b.strip() for b in text.split("\n\n") if b.strip()]
    if not blocks:
        blocks = [text]
    for block in blocks:
        lines = [l for l in block.split("\n") if l.strip()]
        title = lines[0] if lines else "Untitled"
        body = lines[1:] if len(lines) > 1 else []
        _add_slide_with_text(prs, theme, title, body)
    prs.save(output_path)


def _looks_like_heading(paragraph) -> bool:
    """Heuristic: is this docx paragraph likely meant as a section heading?
    True Word heading styles always count. As a fallback (since many real documents never
    use proper Word styles), a short, fully-bolded line with no trailing period is also
    treated as a heading."""
    style_name = paragraph.style.name.lower()
    if style_name.startswith("heading") or style_name == "title":
        return True
    text = paragraph.text.strip()
    if not text or len(text) > 60 or text.endswith((".", ",", ";")):
        return False
    runs = [r for r in paragraph.runs if r.text.strip()]
    return bool(runs) and all(r.bold for r in runs)


def _extract_paragraph_images(paragraph, doc) -> list[bytes]:
    """Extract raw image bytes for any images embedded directly in this docx paragraph."""
    from docx.oxml.ns import qn
    images = []
    for run in paragraph.runs:
        for blip in run._element.findall(".//" + qn("a:blip")):
            rId = blip.get(qn("r:embed"))
            if rId and rId in doc.part.rels:
                images.append(doc.part.rels[rId].target_part.blob)
    return images


def pptx_to_docx(path: str, output_path: str) -> None:
    """Convert a .pptx file to .docx: each slide's title becomes a heading, its text becomes
    paragraphs, and any picture on the slide is inserted into the document."""
    from docx import Document
    prs = Presentation(path)
    doc = Document()
    for i, slide in enumerate(prs.slides):
        texts = []
        image_blobs = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    line = "".join(run.text for run in paragraph.runs)
                    if line:
                        texts.append(line)
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                image_blobs.append(shape.image.blob)
        if not texts and not image_blobs:
            continue
        if texts:
            doc.add_heading(texts[0], level=1)
            for line in texts[1:]:
                doc.add_paragraph(line)
        for blob in image_blobs:
            doc.add_picture(io.BytesIO(blob), width=Inches(4))
        if i < len(prs.slides) - 1:
            doc.add_page_break()
    doc.save(output_path)


def docx_to_pptx(path: str, output_path: str, theme_name: str | None = None) -> None:
    """Convert a .docx file to a themed .pptx: headings (or bolded short lines) start new
    slides, normal paragraphs become bullet content, and any embedded image is placed
    on its corresponding slide. theme_name picks a specific theme; omit for a random one."""
    from docx import Document
    import tempfile

    doc = Document(path)
    theme = _pick_theme(theme_name)
    prs = Presentation()

    current_title = None
    current_body: list[str] = []
    current_image: str | None = None

    def flush():
        if current_title is not None:
            _add_slide_with_text(prs, theme, current_title, current_body, current_image)

    for p in doc.paragraphs:
        text = p.text.strip()
        images = _extract_paragraph_images(p, doc)
        image_path = None
        if images:
            tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".png")
            tmp.write(images[0])
            tmp.close()
            image_path = tmp.name

        if not text and not image_path:
            continue

        if text and _looks_like_heading(p):
            flush()
            current_title = text
            current_body = []
            current_image = image_path
        else:
            if current_title is None:
                current_title = text or "Untitled"
            elif text:
                current_body.append(text)
            if image_path and current_image is None:
                current_image = image_path
    flush()

    prs.save(output_path)


def pdf_to_pptx(path: str, output_path: str, theme_name: str | None = None) -> None:
    """Convert a PDF to a themed .pptx: each page's extracted text and first embedded
    image (if any) are placed onto their own slide. theme_name picks a specific theme;
    omit for a random one."""
    from pypdf import PdfReader
    import tempfile

    reader = PdfReader(path)
    theme = _pick_theme(theme_name)
    prs = Presentation()

    for i, page in enumerate(reader.pages):
        text = page.extract_text() or ""
        lines = [l for l in text.split("\n") if l.strip()]
        title = lines[0] if lines else f"Page {i + 1}"
        body = lines[1:] if len(lines) > 1 else []

        image_path = None
        try:
            images = page.images
            if len(images) > 0:
                tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".png")
                images[0].image.save(tmp.name)
                tmp.close()
                image_path = tmp.name
        except Exception:
            image_path = None

        _add_slide_with_text(prs, theme, title, body, image_path)

    prs.save(output_path)


def extract_pptx_notes(path: str) -> list[str]:
    """Extract the speaker notes from every slide in a .pptx file."""
    prs = Presentation(path)
    notes = []
    for slide in prs.slides:
        if slide.has_notes_slide:
            notes.append(slide.notes_slide.notes_text_frame.text)
        else:
            notes.append("")
    return notes


def merge_pptx_files(paths: list[str], output_path: str, theme_name: str | None = None) -> None:
    """Merge multiple .pptx files' content into one themed deck, preserving slide order."""
    theme = _pick_theme(theme_name)
    prs = Presentation()
    for path in paths:
        source = Presentation(path)
        for slide in source.slides:
            texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        line = "".join(run.text for run in paragraph.runs)
                        if line:
                            texts.append(line)
            title = texts[0] if texts else "Untitled"
            body = texts[1:] if len(texts) > 1 else []
            _add_slide_with_text(prs, theme, title, body)
    prs.save(output_path)